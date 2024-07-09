""" Imports: The necessary modules and packages are imported. This includes PyQt6 for the GUI, pptx for handling PowerPoint files, and deep_translator for translation.

TranslatorThread Class: This class handles the translation process in a separate thread to keep the GUI responsive.

progress, finished, and error signals are defined to communicate with the main thread.
The run method performs the translation and updates the progress.
PPTXTranslatorApp Class: This class represents the main application window.

The initUI method sets up the UI components, including labels, buttons, combo boxes, and progress bar.
select_input_file and select_output_file methods handle file selection dialogs.
translate_pptx method initiates the translation process.
update_progress, translation_finished, and translation_error methods handle updates from the translation thread.
Main Entry Point: The application is created and executed. The main window is displayed, and the application runs until the user closes it."""

import sys
from PyQt6.QtWidgets import (QApplication, QWidget, QVBoxLayout, QHBoxLayout, QPushButton,
                             QFileDialog, QComboBox, QLabel, QProgressBar, QMessageBox)
from PyQt6.QtCore import Qt, QThread, pyqtSignal
from pptx import Presentation
from deep_translator import GoogleTranslator

# This class handles the translation process in a separate thread
class TranslatorThread(QThread):
    # Define custom signals for progress, finish, and error
    progress = pyqtSignal(int)
    finished = pyqtSignal()
    error = pyqtSignal(str)

    def __init__(self, input_file, output_file, source, target):
        QThread.__init__(self)
        self.input_file = input_file
        self.output_file = output_file
        self.source = source
        self.target = target

    def run(self):
        try:
            # Initialize the translator
            translator = GoogleTranslator(source=self.source, target=self.target)
            prs = Presentation(self.input_file)
            total_shapes = sum(len(slide.shapes) for slide in prs.slides)
            processed_shapes = 0

            for slide in prs.slides:
                for shape in slide.shapes:
                    try:
                        # Translate shape text if it exists
                        if hasattr(shape, "text"):
                            if shape.text.strip():
                                shape.text = translator.translate(shape.text)
                        elif hasattr(shape, "text_frame"):
                            if shape.text_frame.text.strip():
                                for paragraph in shape.text_frame.paragraphs:
                                    if paragraph.text.strip():
                                        paragraph.text = translator.translate(paragraph.text)
                    except Exception as shape_error:
                        print(f"Error processing shape: {shape_error}")

                    # Emit progress signal
                    processed_shapes += 1
                    self.progress.emit(int(processed_shapes / total_shapes * 100))

            # Save the translated presentation
            prs.save(self.output_file)
            self.finished.emit()
        except Exception as e:
            self.error.emit(str(e))

# This class represents the main application window
class PPTXTranslatorApp(QWidget):
    def __init__(self):
        super().__init__()
        self.initUI()

    def initUI(self):
        # Set up the main window
        self.setWindowTitle('PPTX Translator')
        self.setGeometry(300, 300, 500, 250)

        layout = QVBoxLayout()

        # Input file selection
        input_file_layout = QHBoxLayout()
        self.input_file_label = QLabel('No input file selected')
        self.input_file_button = QPushButton('Select Input PPTX')
        self.input_file_button.clicked.connect(self.select_input_file)
        input_file_layout.addWidget(self.input_file_label)
        input_file_layout.addWidget(self.input_file_button)
        layout.addLayout(input_file_layout)

        # Output file selection
        output_file_layout = QHBoxLayout()
        self.output_file_label = QLabel('No output file selected')
        self.output_file_button = QPushButton('Select Output PPTX')
        self.output_file_button.clicked.connect(self.select_output_file)
        output_file_layout.addWidget(self.output_file_label)
        output_file_layout.addWidget(self.output_file_button)
        layout.addLayout(output_file_layout)

        # Language selection
        lang_layout = QHBoxLayout()
        self.source_lang = QComboBox()
        self.target_lang = QComboBox()
        languages = ['en', 'es', 'fr', 'de', 'it', 'pt', 'ru', 'ja', 'ko', 'zh']
        self.source_lang.addItems(languages)
        self.target_lang.addItems(languages)
        self.source_lang.setCurrentText('en')
        self.target_lang.setCurrentText('es')
        lang_layout.addWidget(QLabel('From:'))
        lang_layout.addWidget(self.source_lang)
        lang_layout.addWidget(QLabel('To:'))
        lang_layout.addWidget(self.target_lang)
        layout.addLayout(lang_layout)

        # Progress bar
        self.progress_bar = QProgressBar()
        layout.addWidget(self.progress_bar)

        # Translate button
        self.translate_button = QPushButton('Translate')
        self.translate_button.clicked.connect(self.translate_pptx)
        layout.addWidget(self.translate_button)

        self.setLayout(layout)

    # Method to select input file
    def select_input_file(self):
        file_name, _ = QFileDialog.getOpenFileName(self, 'Select Input PPTX file', '', 'PPTX files (*.pptx)')
        if file_name:
            self.input_file_label.setText(file_name)

    # Method to select output file
    def select_output_file(self):
        file_name, _ = QFileDialog.getSaveFileName(self, 'Select Output PPTX file', '', 'PPTX files (*.pptx)')
        if file_name:
            self.output_file_label.setText(file_name)

    # Method to start translation
    def translate_pptx(self):
        input_file = self.input_file_label.text()
        output_file = self.output_file_label.text()
        if input_file == 'No input file selected' or output_file == 'No output file selected':
            QMessageBox.warning(self, 'Error', 'Please select both input and output files.')
            return

        source = self.source_lang.currentText()
        target = self.target_lang.currentText()

        self.thread = TranslatorThread(input_file, output_file, source, target)
        self.thread.progress.connect(self.update_progress)
        self.thread.finished.connect(self.translation_finished)
        self.thread.error.connect(self.translation_error)

        self.translate_button.setEnabled(False)
        self.progress_bar.setValue(0)
        self.thread.start()

    # Method to update progress bar
    def update_progress(self, value):
        self.progress_bar.setValue(value)

    # Method called when translation is finished
    def translation_finished(self):
        self.translate_button.setEnabled(True)
        QMessageBox.information(self, 'Success', f"Translation completed. File saved as: {self.output_file_label.text()}")

    # Method called when an error occurs during translation
    def translation_error(self, error_message):
        self.translate_button.setEnabled(True)
        QMessageBox.critical(self, 'Error', f"An error occurred during translation: {error_message}")

# Main entry point of the application
if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = PPTXTranslatorApp()
    ex.show()
    sys.exit(app.exec())
